"""Turn a source file into the markdown-ish text the LLM will read.

Text formats are handled with the standard library. Binary formats (PDF,
DOCX, XLSX, PPTX) need a parser library each, and every one of them is an
optional extra — an unreadable format raises `ParseError` with the exact pip
command rather than silently ingesting an empty document.

Extraction preserves structure where it can (headings, lists, tables),
because the ingest prompt asks the model to keep tables and schemas verbatim
and it cannot preserve what the parser already flattened.
"""

from __future__ import annotations

import csv
import html
import io
import json
import re
from html.parser import HTMLParser
from pathlib import Path

from .errors import ParseError

TEXT_EXTENSIONS = {
    ".md", ".markdown", ".txt", ".org", ".rst", ".text", ".log",
    ".py", ".js", ".ts", ".tsx", ".jsx", ".rs", ".go", ".java", ".rb", ".sh",
    ".sql", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".xml",
}
JSON_EXTENSIONS = {".json", ".jsonl", ".ndjson"}
TABLE_EXTENSIONS = {".csv", ".tsv"}
HTML_EXTENSIONS = {".html", ".htm"}

BINARY_HINTS = {
    ".pdf": ("pypdf", "pip install 'llmwiki[docs]'"),
    ".docx": ("python-docx", "pip install 'llmwiki[docs]'"),
    ".xlsx": ("openpyxl", "pip install 'llmwiki[docs]'"),
    ".xlsm": ("openpyxl", "pip install 'llmwiki[docs]'"),
    ".pptx": ("python-pptx", "pip install 'llmwiki[docs]'"),
}


def supported_extensions() -> set[str]:
    return (
        TEXT_EXTENSIONS
        | JSON_EXTENSIONS
        | TABLE_EXTENSIONS
        | HTML_EXTENSIONS
        | set(BINARY_HINTS)
    )


def extract_text(path: Path) -> str:
    """Read `path` as text. Raises `ParseError` if it cannot be read."""
    suffix = path.suffix.lower()
    if suffix in TEXT_EXTENSIONS:
        return _read_text(path)
    if suffix in JSON_EXTENSIONS:
        return _read_json(path)
    if suffix in TABLE_EXTENSIONS:
        return _read_table(path, delimiter="\t" if suffix == ".tsv" else ",")
    if suffix in HTML_EXTENSIONS:
        return html_to_markdown(_read_text(path))
    if suffix == ".pdf":
        return _read_pdf(path)
    if suffix == ".docx":
        return _read_docx(path)
    if suffix in {".xlsx", ".xlsm"}:
        return _read_xlsx(path)
    if suffix == ".pptx":
        return _read_pptx(path)

    # Unknown extension: try UTF-8 and accept it if it looks like text.
    try:
        text = path.read_text(encoding="utf-8")
    except (OSError, UnicodeDecodeError) as exc:
        raise ParseError(
            f"cannot read {path.name}: unsupported format '{suffix or 'no extension'}'"
        ) from exc
    if "\x00" in text:
        raise ParseError(f"cannot read {path.name}: looks like binary data")
    return text


def _read_text(path: Path) -> str:
    for encoding in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
        except OSError as exc:
            raise ParseError(f"cannot read {path.name}: {exc}") from exc
    raise ParseError(f"cannot decode {path.name} as text")


def _read_json(path: Path) -> str:
    raw = _read_text(path)
    try:
        parsed = json.loads(raw)
    except json.JSONDecodeError:
        return f"```json\n{raw}\n```"  # jsonl or malformed: hand it over as-is
    return f"```json\n{json.dumps(parsed, indent=2, ensure_ascii=False)}\n```"


def _read_table(path: Path, delimiter: str) -> str:
    """CSV to a markdown table — the shape the model reasons about best."""
    raw = _read_text(path)
    rows = list(csv.reader(io.StringIO(raw), delimiter=delimiter))
    if not rows:
        return ""
    header, *body = rows
    lines = [
        "| " + " | ".join(_escape_cell(c) for c in header) + " |",
        "| " + " | ".join("---" for _ in header) + " |",
    ]
    for row in body:
        padded = row + [""] * (len(header) - len(row))
        lines.append("| " + " | ".join(_escape_cell(c) for c in padded[: len(header)]) + " |")
    return "\n".join(lines)


def _escape_cell(value: str) -> str:
    return value.replace("|", "\\|").replace("\n", " ").strip()


class _HTMLToMarkdown(HTMLParser):
    """Enough of Readability+Turndown to keep headings, lists, and links."""

    _SKIP = {"script", "style", "nav", "header", "footer", "aside", "noscript"}

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.parts: list[str] = []
        self._skip_depth = 0
        self._list_stack: list[str] = []
        self._href: str | None = None

    def handle_starttag(self, tag, attrs):
        if tag in self._SKIP:
            self._skip_depth += 1
            return
        if self._skip_depth:
            return
        if tag in {"h1", "h2", "h3", "h4", "h5", "h6"}:
            self.parts.append(f"\n\n{'#' * int(tag[1])} ")
        elif tag in {"p", "div", "section", "article", "br", "tr"}:
            self.parts.append("\n\n" if tag != "br" else "\n")
        elif tag in {"ul", "ol"}:
            self._list_stack.append(tag)
            self.parts.append("\n")
        elif tag == "li":
            marker = "1." if (self._list_stack and self._list_stack[-1] == "ol") else "-"
            self.parts.append(f"\n{'  ' * max(0, len(self._list_stack) - 1)}{marker} ")
        elif tag == "a":
            self._href = dict(attrs).get("href")
            self.parts.append("[")
        elif tag in {"strong", "b"}:
            self.parts.append("**")
        elif tag in {"em", "i"}:
            self.parts.append("*")
        elif tag in {"code", "pre"}:
            self.parts.append("`" if tag == "code" else "\n\n```\n")

    def handle_endtag(self, tag):
        if tag in self._SKIP:
            self._skip_depth = max(0, self._skip_depth - 1)
            return
        if self._skip_depth:
            return
        if tag in {"ul", "ol"} and self._list_stack:
            self._list_stack.pop()
        elif tag == "a":
            self.parts.append(f"]({self._href})" if self._href else "]")
            self._href = None
        elif tag in {"strong", "b"}:
            self.parts.append("**")
        elif tag in {"em", "i"}:
            self.parts.append("*")
        elif tag == "code":
            self.parts.append("`")
        elif tag == "pre":
            self.parts.append("\n```\n")

    def handle_data(self, data):
        if self._skip_depth:
            return
        self.parts.append(data)

    def markdown(self) -> str:
        text = "".join(self.parts)
        text = re.sub(r"[ \t]+", " ", text)
        text = re.sub(r"\n{3,}", "\n\n", text)
        return html.unescape(text).strip()


def html_to_markdown(raw_html: str) -> str:
    parser = _HTMLToMarkdown()
    parser.feed(raw_html)
    parser.close()
    return parser.markdown()


def _missing(path: Path, module: str, hint: str) -> ParseError:
    return ParseError(f"reading {path.name} needs the '{module}' package — {hint}")


def _read_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise _missing(path, *BINARY_HINTS[".pdf"]) from exc
    try:
        reader = PdfReader(str(path))
        pages = []
        for number, page in enumerate(reader.pages, start=1):
            text = (page.extract_text() or "").strip()
            if text:
                pages.append(f"## Page {number}\n\n{text}")
        return "\n\n".join(pages)
    except Exception as exc:
        raise ParseError(f"could not extract text from {path.name}: {exc}") from exc


def _read_docx(path: Path) -> str:
    try:
        import docx
    except ImportError as exc:
        raise _missing(path, *BINARY_HINTS[".docx"]) from exc
    try:
        document = docx.Document(str(path))
    except Exception as exc:
        raise ParseError(f"could not open {path.name}: {exc}") from exc

    lines: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style = (paragraph.style.name or "").lower()
        if style.startswith("heading"):
            level = "".join(ch for ch in style if ch.isdigit()) or "1"
            lines.append(f"{'#' * min(6, int(level))} {text}")
        elif style.startswith("list"):
            lines.append(f"- {text}")
        else:
            lines.append(text)
    for table in document.tables:
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        if not rows:
            continue
        lines.append("")
        lines.append("| " + " | ".join(_escape_cell(c) for c in rows[0]) + " |")
        lines.append("| " + " | ".join("---" for _ in rows[0]) + " |")
        for row in rows[1:]:
            lines.append("| " + " | ".join(_escape_cell(c) for c in row) + " |")
    return "\n\n".join(lines)


def _read_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise _missing(path, *BINARY_HINTS[".xlsx"]) from exc
    try:
        workbook = load_workbook(str(path), read_only=True, data_only=True)
    except Exception as exc:
        raise ParseError(f"could not open {path.name}: {exc}") from exc

    out: list[str] = []
    for sheet in workbook.worksheets:
        rows = [
            ["" if cell is None else str(cell) for cell in row]
            for row in sheet.iter_rows(values_only=True)
        ]
        rows = [row for row in rows if any(cell.strip() for cell in row)]
        if not rows:
            continue
        out.append(f"## {sheet.title}")
        header = rows[0]
        out.append("| " + " | ".join(_escape_cell(c) for c in header) + " |")
        out.append("| " + " | ".join("---" for _ in header) + " |")
        for row in rows[1:]:
            padded = row + [""] * (len(header) - len(row))
            out.append("| " + " | ".join(_escape_cell(c) for c in padded[: len(header)]) + " |")
        out.append("")
    return "\n".join(out)


def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise _missing(path, *BINARY_HINTS[".pptx"]) from exc
    try:
        presentation = Presentation(str(path))
    except Exception as exc:
        raise ParseError(f"could not open {path.name}: {exc}") from exc

    out: list[str] = []
    for number, slide in enumerate(presentation.slides, start=1):
        out.append(f"## Slide {number}")
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = "".join(run.text for run in paragraph.runs).strip()
                if text:
                    out.append(f"- {text}" if paragraph.level else text)
        out.append("")
    return "\n".join(out)
